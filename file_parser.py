"""
File parser - Extract text content and embedded images from various file types.
"""
import os
import tempfile
import chardet
from typing import Optional

# Max characters to extract from a file
MAX_TEXT_LENGTH = 4000

# File extensions by category
TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".log", ".ini", ".yaml", ".yml",
    ".toml", ".cfg", ".conf", ".env", ".gitignore", ".dockerfile",
    ".rst", ".tex", ".bib",
}

CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm", ".css", ".scss",
    ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".rb",
    ".php", ".swift", ".kt", ".scala", ".lua", ".r", ".m", ".sql",
    ".sh", ".bat", ".ps1", ".makefile", ".cmake",
    ".vue", ".svelte", ".astro",
}

IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".ico",
}

DOCUMENT_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
}


def get_file_category(file_path: str) -> str:
    """Determine the file category based on extension."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in TEXT_EXTENSIONS:
        return "text"
    if ext in CODE_EXTENSIONS:
        return "code"
    if ext in DOCUMENT_EXTENSIONS:
        return "document"
    return "unknown"


def parse_file(file_path: str) -> tuple[Optional[str], str]:
    """
    Parse a file and return (text_content, category).
    For images, text_content is None (handled by vision model).
    """
    category = get_file_category(file_path)

    if category == "image":
        return None, "image"

    try:
        if category in ("text", "code"):
            return _read_text_file(file_path), category
        elif category == "document":
            return _parse_document(file_path), category
        else:
            # Unknown: try to read as text
            return _read_text_file(file_path), "text"
    except Exception as e:
        print(f"[Parser] Failed to parse {file_path}: {e}")
        return None, "error"


def _read_text_file(file_path: str) -> Optional[str]:
    """Read a text file with encoding detection."""
    try:
        with open(file_path, "rb") as f:
            raw = f.read(100000)  # Read up to 100KB for detection

        if not raw:
            return None

        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"

        try:
            text = raw.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            text = raw.decode("utf-8", errors="replace")

        return text[:MAX_TEXT_LENGTH] if text.strip() else None
    except Exception:
        return None


def _parse_document(file_path: str) -> Optional[str]:
    """Parse document files (PDF, DOCX, XLSX, PPTX) - text extraction only."""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return _parse_pdf(file_path)
        elif ext == ".docx":
            return _parse_docx(file_path)
        elif ext == ".xlsx":
            return _parse_xlsx(file_path)
        elif ext == ".pptx":
            return _parse_pptx(file_path)
    except Exception as e:
        print(f"[Parser] Error parsing document {file_path}: {e}")
    return None


def _parse_pdf(file_path: str) -> Optional[str]:
    """Extract text from PDF using pypdf."""
    import pypdf

    text_parts = []
    try:
        reader = pypdf.PdfReader(file_path)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
    except Exception as e:
        print(f"[Parser] PDF error: {e}")
        return None

    text = "\n".join(text_parts)
    return text[:MAX_TEXT_LENGTH] if text.strip() else None


def _parse_docx(file_path: str) -> Optional[str]:
    """Extract text from DOCX."""
    from docx import Document

    doc = Document(file_path)
    text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n".join(text_parts)
    return text[:MAX_TEXT_LENGTH] if text.strip() else None


def _parse_xlsx(file_path: str) -> Optional[str]:
    """Extract text from XLSX."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text_parts.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                text_parts.append(" | ".join(values))
    wb.close()

    text = "\n".join(text_parts)
    return text[:MAX_TEXT_LENGTH] if text.strip() else None


def _parse_pptx(file_path: str) -> Optional[str]:
    """Extract text from PPTX."""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)

    text = "\n".join(text_parts)
    return text[:MAX_TEXT_LENGTH] if text.strip() else None


# ──────────────── Document → Images (for Vision LLM) ────────────────


def get_pdf_page_images(file_path: str, max_pages: int = 5, dpi: int = 200) -> list[str]:
    """
    Convert PDF pages to images using pymupdf.
    Returns list of temporary image file paths.
    """
    import fitz  # pymupdf

    image_paths = []
    try:
        doc = fitz.open(file_path)
        num_pages = min(len(doc), max_pages)
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)

        for page_num in range(num_pages):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=matrix)

            # Save to temp file
            tmp_path = os.path.join(
                tempfile.gettempdir(),
                f"fileguessr_pdf_p{page_num}_{os.path.basename(file_path)}.png"
            )
            pix.save(tmp_path)
            image_paths.append(tmp_path)

        doc.close()
    except Exception as e:
        print(f"[Parser] Error converting PDF to images: {e}")

    return image_paths


def get_docx_images(file_path: str) -> list[str]:
    """
    Extract embedded images from a DOCX file.
    Returns list of temporary image file paths.
    """
    from docx import Document

    image_paths = []
    try:
        doc = Document(file_path)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    # Determine extension from content type
                    content_type = rel.target_part.content_type
                    ext = ".png"
                    if "jpeg" in content_type or "jpg" in content_type:
                        ext = ".jpg"
                    elif "gif" in content_type:
                        ext = ".gif"
                    elif "bmp" in content_type:
                        ext = ".bmp"

                    tmp_path = os.path.join(
                        tempfile.gettempdir(),
                        f"fileguessr_docx_{len(image_paths)}_{os.path.basename(file_path)}{ext}"
                    )
                    with open(tmp_path, "wb") as f:
                        f.write(image_data)
                    image_paths.append(tmp_path)
                except Exception as e:
                    print(f"[Parser] Error extracting DOCX image: {e}")
    except Exception as e:
        print(f"[Parser] Error reading DOCX images: {e}")

    return image_paths


def get_pptx_images(file_path: str) -> list[str]:
    """
    Extract embedded images from a PPTX file.
    Returns list of temporary image file paths.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    image_paths = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        tmp_path = os.path.join(
                            tempfile.gettempdir(),
                            f"fileguessr_pptx_{len(image_paths)}_{os.path.basename(file_path)}.{ext}"
                        )
                        with open(tmp_path, "wb") as f:
                            f.write(image.blob)
                        image_paths.append(tmp_path)
                    except Exception as e:
                        print(f"[Parser] Error extracting PPTX image: {e}")
    except Exception as e:
        print(f"[Parser] Error reading PPTX images: {e}")

    return image_paths


def get_xlsx_images(file_path: str) -> list[str]:
    """
    Extract embedded images from an XLSX file.
    Returns list of temporary image file paths.
    """
    from openpyxl import load_workbook

    image_paths = []
    try:
        wb = load_workbook(file_path)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for image in ws._images:
                try:
                    tmp_path = os.path.join(
                        tempfile.gettempdir(),
                        f"fileguessr_xlsx_{len(image_paths)}_{os.path.basename(file_path)}.png"
                    )
                    with open(tmp_path, "wb") as f:
                        f.write(image._data())
                    image_paths.append(tmp_path)
                except Exception as e:
                    print(f"[Parser] Error extracting XLSX image: {e}")
        wb.close()
    except Exception as e:
        print(f"[Parser] Error reading XLSX images: {e}")

    return image_paths


def get_document_images(file_path: str) -> list[str]:
    """
    Extract images from any document type.
    For PDF: renders each page as an image.
    For DOCX/PPTX/XLSX: extracts embedded images.
    Returns list of temporary image file paths.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return get_pdf_page_images(file_path)
    elif ext == ".docx":
        return get_docx_images(file_path)
    elif ext == ".pptx":
        return get_pptx_images(file_path)
    elif ext == ".xlsx":
        return get_xlsx_images(file_path)
    return []


def cleanup_temp_images(image_paths: list[str]):
    """Clean up temporary image files."""
    for path in image_paths:
        try:
            if os.path.exists(path) and path.startswith(tempfile.gettempdir()):
                os.remove(path)
        except Exception:
            pass
